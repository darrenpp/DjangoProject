from __future__ import annotations

import shutil
from datetime import datetime
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
PRESENTATIONS_DIR = Path(r"C:\Users\darre\Documents\ProjectApps\databasedocuments\Presentations")
SOURCE_PPTX = PRESENTATIONS_DIR / "NDOH_Regulatory_Platform_Presentation_for_Registrar_20260604_v2.pptx"
CORRECTED_PPTX = SOURCE_PPTX
HANDOUT_PDF = PRESENTATIONS_DIR / "NDOH_Regulatory_Platform_Attendee_Handout.pdf"
HANDOUT_2UP_PDF = PRESENTATIONS_DIR / "NDOH_Regulatory_Platform_Attendee_Handout_2up_Print.pdf"

SCREENSHOTS = ROOT / "docs" / "presentation" / "assets" / "screenshots"
EMBLEM = ROOT / "static" / "img" / "National_emblem_of_Papua_New_Guinea_(variant).svg.png"
WORK = ROOT / "generated_reports" / "twg_presentation_20260609" / "assets"


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        Path("C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf"),
        Path("C:/Windows/Fonts/calibrib.ttf" if bold else "C:/Windows/Fonts/calibri.ttf"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size)
    return ImageFont.load_default()


def backup(path: Path) -> Path | None:
    if not path.exists():
        return None
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = path.with_suffix(path.suffix + f".bak_{stamp}")
    shutil.copy2(path, backup_path)
    return backup_path


def remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def picture_shapes(slide):
    return [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]


def save_contained(src: Path, dest: Path, size: tuple[int, int], background: str = "#f8fafc") -> Path:
    image = Image.open(src).convert("RGB")
    contained = ImageOps.contain(image, size, method=Image.Resampling.LANCZOS)
    canvas = Image.new("RGB", size, background)
    x = (size[0] - contained.width) // 2
    y = (size[1] - contained.height) // 2
    canvas.paste(contained, (x, y))
    canvas.save(dest)
    return dest


def make_platform_montage() -> Path:
    WORK.mkdir(parents=True, exist_ok=True)
    output = WORK / "current_platform_montage_20260610.png"
    cells = [
        ("System Admin", SCREENSHOTS / "overall_dashboard.png"),
        ("Nursing Council", SCREENSHOTS / "nursing_council_dashboard.png"),
        ("Medical Board", SCREENSHOTS / "medical_board_dashboard.png"),
        ("Records Hub", SCREENSHOTS / "records_hub.png"),
    ]
    width, height = 1800, 1330
    canvas = Image.new("RGB", (width, height), "#eef6f8")
    draw = ImageDraw.Draw(canvas)
    draw.rounded_rectangle((20, 20, width - 20, height - 20), radius=26, fill="#ffffff", outline="#cbd5e1", width=3)
    draw.text((54, 40), "Current Platform Interfaces", font=font(44, True), fill="#0a2b45")
    draw.text((54, 94), "Refreshed from the local platform design on 10 June 2026.", font=font(24), fill="#476173")

    cell_w, cell_h = 825, 520
    positions = [(54, 170), (920, 170), (54, 745), (920, 745)]
    for (label, path), (x, y) in zip(cells, positions):
        draw.rounded_rectangle((x, y, x + cell_w, y + cell_h), radius=18, fill="#f8fafc", outline="#d7e4ea", width=2)
        if path.exists():
            shot = Image.open(path).convert("RGB")
            shot = ImageOps.contain(shot, (cell_w - 34, cell_h - 88), method=Image.Resampling.LANCZOS)
            canvas.paste(shot, (x + 17 + (cell_w - 34 - shot.width) // 2, y + 18))
        draw.rounded_rectangle((x + 18, y + cell_h - 58, x + cell_w - 18, y + cell_h - 18), radius=12, fill="#0f766e")
        draw.text((x + 38, y + cell_h - 48), label, font=font(24, True), fill="#ffffff")
    canvas.save(output)
    return output


def patch_deck() -> Path:
    WORK.mkdir(parents=True, exist_ok=True)
    montage = make_platform_montage()
    repo_view = save_contained(
        SCREENSHOTS / "documents_search.png",
        WORK / "current_documents_repository_20260610.png",
        (1800, 1465),
    )

    prs = Presentation(SOURCE_PPTX)

    first_slide = prs.slides[0]
    first_pictures = picture_shapes(first_slide)
    if len(first_pictures) >= 2 and EMBLEM.exists():
        target = first_pictures[-1]
        left, top, width, height = target.left, target.top, target.width, target.height
        remove_shape(target)
        first_slide.shapes.add_picture(str(EMBLEM), left, top, width=width, height=height)

    slide_8 = prs.slides[7]
    slide_8_pictures = picture_shapes(slide_8)
    for target in slide_8_pictures:
        left, top, width, height = target.left, target.top, target.width, target.height
        remove_shape(target)
        slide_8.shapes.add_picture(str(montage), left, top, width=width, height=height)

    slide_10 = prs.slides[9]
    slide_10_pictures = picture_shapes(slide_10)
    for target in slide_10_pictures:
        left, top, width, height = target.left, target.top, target.width, target.height
        remove_shape(target)
        slide_10.shapes.add_picture(str(repo_view), left, top, width=width, height=height)

    temp_pptx = WORK / "NDOH_Regulatory_Platform_Presentation_for_Registrar_20260604_v2_final.tmp.pptx"
    prs.save(temp_pptx)
    shutil.move(str(temp_pptx), str(CORRECTED_PPTX))
    return CORRECTED_PPTX


def main() -> None:
    backup(HANDOUT_PDF)
    backup(HANDOUT_2UP_PDF)
    corrected = patch_deck()
    print(f"CORRECTED_PPTX={corrected}")
    print(f"HANDOUT_PDF={HANDOUT_PDF}")
    print(f"HANDOUT_2UP_PDF={HANDOUT_2UP_PDF}")


if __name__ == "__main__":
    main()
