from __future__ import annotations

import shutil
from datetime import datetime
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "generated_reports" / "twg_presentation_20260609"
SCREENSHOTS = ROOT / "docs" / "presentation" / "assets" / "screenshots"
DIAGRAMS = ROOT / "docs" / "presentation" / "assets" / "diagrams"
DATE_TEXT = "9 June 2026"
ANNEX_LABEL = "Annex D"

NAVY = PptxRGBColor(10, 43, 69)
TEAL = PptxRGBColor(20, 90, 99)
GREEN = PptxRGBColor(22, 116, 92)
GOLD = PptxRGBColor(176, 126, 36)
LIGHT = PptxRGBColor(238, 246, 248)
DARK_TEXT = PptxRGBColor(11, 31, 47)


SLIDES = [
    {
        "title": "TWG Presentation Focus",
        "subtitle": "PNG Medical Board & Nursing Council Regulatory Agencies Platform",
        "bullets": [
            "Purpose, current status, readiness, and controlled testing pathway.",
            "Clear separation of Medical Board and Nursing Council workspaces.",
            "ICT support needed before production deployment.",
        ],
        "notes": "Open by positioning the system as a regulatory operations platform, not a general website. The main message is that the foundation is in place for controlled testing, while production depends on ICT review and formal sign-off.",
        "image": "public_home.png",
    },
    {
        "title": "Purpose Of The Platform",
        "bullets": [
            "Provide one controlled digital foundation for professional registration, licensing, applications, records, documents, receipts, data quality, and reporting.",
            "Support both regulatory agencies while preserving workspace separation and privacy boundaries.",
            "Move from spreadsheet/manual processing toward auditable, role-based regulatory workflows.",
        ],
        "notes": "Explain that the goal is not to replace registrar authority. The platform supports registrar decisions by organising data, workflow, evidence, and reporting.",
    },
    {
        "title": "Current Status",
        "bullets": [
            "Core platform foundation is implemented and running locally.",
            "Controlled testing and UAT can begin with authorised staff.",
            "Production use is not yet recommended until ICT hosting, security, backup, UAT, training, and readiness gates are completed.",
        ],
        "notes": "Use careful language: ready for controlled testing, not yet approved for live national production.",
        "image": "production_readiness.png",
    },
    {
        "title": "Functional Coverage",
        "bullets": [
            "Registration and licence pathways for Nursing Council and Medical Board users.",
            "Application review, checklist, payment verification, document issue, and status history.",
            "Records hub, public-safe search, reports, financial views, maps, notifications, documents, and complaints/discipline registers.",
        ],
        "notes": "Summarise the breadth without diving into every screen. The purpose is to show that the platform is an integrated regulatory workflow system.",
        "image": "overall_dashboard.png",
    },
    {
        "title": "Separated Workspaces",
        "bullets": [
            "Nursing Council workflows cover nurses, midwives, nurse aides, graduands, ATP/practising licence, and nursing-specific reference data.",
            "Medical Board workflows cover doctors, specialists, CHWs, and medical board-specific records and references.",
            "Cross-office access should only occur where formally authorised.",
        ],
        "notes": "Stress that privacy and regulatory authority are preserved. Nursing Council users should not automatically access private Medical Board records, and Medical Board users should not automatically access private Nursing Council records.",
        "image": "medical_board_dashboard.png",
    },
    {
        "title": "Registration, Licensing, And Applications",
        "bullets": [
            "Supports provisional, full licence, ATP/practising, renewal, graduand/student, CHW, doctor, and nurse aide pathways.",
            "Application records can carry status, documents, checklist review, payment verification, and registrar decision history.",
            "Professional users can access their own profile, applications, receipts, and documents.",
        ],
        "notes": "Link this to the user journey: applicant submits, staff review, registrar decides, status and audit trail are preserved.",
        "image": "nursing_forms.png",
    },
    {
        "title": "Records, Documents, And Receipts",
        "bullets": [
            "Records Hub supports controlled registry and reference-table management.",
            "Document repository supports metadata, versions, approval actions, OCR/search readiness, and audit history.",
            "Receipt linking and finance views support payment evidence and review routing.",
        ],
        "notes": "Explain that records and documents are treated as evidence, not casual uploads. Receipts and documents need controlled review before they support operational decisions.",
        "image": "records_hub.png",
    },
    {
        "title": "Data Quality And Reporting",
        "bullets": [
            "Imported records are staged, checked, cleansed, and reviewed before promotion into operational records.",
            "Duplicate review, missing-data review, source traceability, and analytics snapshots support data-quality improvement.",
            "Reports and dashboards support registrar oversight, finance, workforce planning, and NHWA-aligned reporting.",
        ],
        "notes": "This is a key valuation and governance point. The platform is not just storing data; it provides a controlled process for improving and defending data quality.",
        "image": "duplicate_review_queue.png",
    },
    {
        "title": "Role-Based Access, Privacy, And Audit",
        "bullets": [
            "Role-based dashboards separate public, professional, registrar/staff, reviewer, finance, data-quality, and system-admin functions.",
            "Staff accounts require approval before login and operational access can be requested separately.",
            "Audit evidence includes login/security events, application status history, document audit events, and regulatory logs.",
        ],
        "notes": "This slide should reassure the TWG that privacy and audit have been considered, while also making clear that formal security review is still required.",
        "image": "staff_notifications.png",
    },
    {
        "title": "Data Governance Process",
        "bullets": [
            "Import source data.",
            "Stage and validate rows.",
            "Clean names, roles, institutions, facilities, and registration details.",
            "Review duplicates, missing data, and source conflicts.",
            "Promote only authorised records into operational registry use.",
        ],
        "notes": "The strongest governance message is that imported spreadsheet rows do not automatically become legal registry records. Promotion must be controlled and authorised.",
        "image": "nhwa_workbooks.png",
    },
    {
        "title": "Controlled Testing And UAT Readiness",
        "bullets": [
            "225 automated Django tests passed during the evidence-pack preparation.",
            "Mobile helper confirmed local Android testing settings, 22 enabled mobile form schemas, and a passing health endpoint.",
            "UAT should cover registrar, finance, data quality, professional user, public register, records, documents, reports, and mobile sync workflows.",
        ],
        "notes": "Use this as evidence of readiness for controlled testing. Do not present automated tests as a replacement for user acceptance testing.",
        "image": "nurse_portal.png",
    },
    {
        "title": "Support Required From NDOH ICT",
        "bullets": [
            "On-prem hosting assessment and approved server environment.",
            "Security review, vulnerability scan, and independent penetration test.",
            "Backup and restore confirmation, production email, HTTPS/domain, DNS, and deployment planning.",
            "Operational monitoring, support ownership, incident response, and change-control process.",
        ],
        "notes": "Make this practical. The TWG should leave with a clear ICT action list, not just a technical demonstration.",
    },
    {
        "title": "Possible DICT Cloud Pathway",
        "bullets": [
            "Use on-prem or local controlled testing first to validate workflow, data, privacy, and UAT requirements.",
            "Prepare architecture, security, data classification, backup, and access-control evidence for DICT assessment.",
            "Consider DICT cloud only after risk, hosting, compliance, and support responsibilities are documented.",
        ],
        "notes": "Frame DICT cloud as a future assessed pathway, not an immediate shortcut. Production hosting needs a governance decision.",
    },
    {
        "title": "Production Gate",
        "bullets": [
            "Foundation is in place and ready for controlled testing.",
            "Full production use should proceed only after ICT review, UAT sign-off, security checks, backup/restore confirmation, staff training, and production-readiness approval.",
            "TWG decision requested: approve controlled testing pathway and assign owners for ICT, UAT, security, hosting, and training workstreams.",
        ],
        "notes": "Close with a decision request. The desired outcome is permission to proceed with controlled testing and a named owner for each readiness gate.",
    },
]


def reset_output() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    (OUT / "assets").mkdir(exist_ok=True)


def screenshot(name: str) -> Path | None:
    if not name:
        return None
    path = SCREENSHOTS / name
    return path if path.is_file() else None


def set_fill(shape, color: PptxRGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    left = PptxInches(0.55)
    top = PptxInches(0.35)
    width = PptxInches(10.8)
    height = PptxInches(0.7)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.bold = True
    p.font.size = PptxPt(29)
    p.font.color.rgb = NAVY
    if subtitle:
        sub = slide.shapes.add_textbox(left, PptxInches(1.02), width, PptxInches(0.42))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Aptos"
        sp.font.size = PptxPt(14)
        sp.font.color.rgb = TEAL


def add_annex_label(slide) -> None:
    label = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        PptxInches(11.55),
        PptxInches(0.36),
        PptxInches(1.2),
        PptxInches(0.36),
    )
    label.fill.solid()
    label.fill.fore_color.rgb = PptxRGBColor(255, 255, 255)
    label.line.color.rgb = GOLD
    label.text_frame.clear()
    paragraph = label.text_frame.paragraphs[0]
    paragraph.text = ANNEX_LABEL
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = "Aptos"
    paragraph.font.bold = True
    paragraph.font.size = PptxPt(12)
    paragraph.font.color.rgb = NAVY


def add_footer(slide, slide_no: int) -> None:
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        PptxInches(0),
        PptxInches(7.18),
        PptxInches(13.333),
        PptxInches(0.32),
    )
    set_fill(line, NAVY)
    box = slide.shapes.add_textbox(PptxInches(0.55), PptxInches(7.22), PptxInches(10.5), PptxInches(0.22))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "PNG Medical Board & Nursing Council Regulatory Agencies Platform | TWG briefing"
    p.font.name = "Aptos"
    p.font.size = PptxPt(8)
    p.font.color.rgb = PptxRGBColor(255, 255, 255)
    no = slide.shapes.add_textbox(PptxInches(12.25), PptxInches(7.22), PptxInches(0.5), PptxInches(0.22))
    ntf = no.text_frame
    ntf.clear()
    np = ntf.paragraphs[0]
    np.text = str(slide_no)
    np.alignment = PP_ALIGN.RIGHT
    np.font.name = "Aptos"
    np.font.size = PptxPt(8)
    np.font.color.rgb = PptxRGBColor(255, 255, 255)


def add_bullets(slide, bullets: list[str], has_image: bool) -> None:
    left = PptxInches(0.75)
    top = PptxInches(1.65)
    width = PptxInches(6.7 if has_image else 11.6)
    height = PptxInches(4.75)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = PptxPt(18 if len(bullets) <= 3 else 16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = PptxPt(10)


def add_image(slide, image_name: str | None) -> None:
    if not image_name:
        return
    path = screenshot(image_name)
    if not path:
        return
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        PptxInches(8.0),
        PptxInches(1.55),
        PptxInches(4.75),
        PptxInches(4.75),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT
    panel.line.color.rgb = PptxRGBColor(204, 220, 225)
    slide.shapes.add_picture(str(path), PptxInches(8.18), PptxInches(1.78), width=PptxInches(4.38))


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    blank = prs.slide_layouts[6]
    for idx, item in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = PptxRGBColor(249, 252, 253)
        accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptxInches(0), PptxInches(0), PptxInches(0.18), PptxInches(7.5))
        set_fill(accent, TEAL if idx % 3 else GREEN)
        add_title(slide, item["title"], item.get("subtitle"))
        add_annex_label(slide)
        add_bullets(slide, item["bullets"], bool(item.get("image")))
        add_image(slide, item.get("image"))
        add_footer(slide, idx)

    path = OUT / "TWG_Platform_Readiness_Presentation_20260609.pptx"
    try:
        prs.save(path)
    except PermissionError:
        path = OUT / "TWG_Platform_Readiness_Presentation_20260610_corrected.pptx"
        prs.save(path)
    return path


def build_markdown() -> Path:
    lines = [
        "# TWG Platform Readiness Presentation",
        "",
        f"Prepared: {DATE_TEXT}",
        "",
        "## Core Message",
        "",
        "The platform foundation is in place and ready for controlled testing. Full production use should only proceed after NDOH ICT review, UAT sign-off, security checks, backup/restore confirmation, staff training, and production-readiness approval.",
        "",
    ]
    for idx, item in enumerate(SLIDES, 1):
        lines.extend([f"## Slide {idx}: {item['title']}", ""])
        if item.get("subtitle"):
            lines.extend([str(item["subtitle"]), ""])
        lines.extend(f"- {bullet}" for bullet in item["bullets"])
        lines.extend(["", f"Speaker note: {item['notes']}", ""])
    path = OUT / "TWG_Platform_Readiness_Presentation_20260609.md"
    write(path, "\n".join(lines))
    return path


def build_docx() -> Path:
    document = Document()
    section = document.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    annex = document.add_paragraph(ANNEX_LABEL)
    annex.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    annex.runs[0].bold = True
    annex.runs[0].font.size = Pt(10)
    annex.runs[0].font.color.rgb = RGBColor(176, 126, 36)

    title = document.add_heading("TWG Platform Readiness Presentation Brief", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle = document.add_paragraph("PNG Medical Board & Nursing Council Regulatory Agencies Platform")
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.runs[0].font.size = Pt(12)
    subtitle.runs[0].font.color.rgb = RGBColor(20, 90, 99)
    date = document.add_paragraph(f"Prepared: {DATE_TEXT}")
    date.alignment = WD_ALIGN_PARAGRAPH.CENTER

    document.add_heading("Core Message", level=1)
    document.add_paragraph(
        "The platform foundation is in place and ready for controlled testing. Full production use should only proceed after NDOH ICT review, UAT sign-off, security checks, backup/restore confirmation, staff training, and production-readiness approval."
    )

    document.add_heading("TWG Outcomes Requested", level=1)
    for item in [
        "Agree to proceed with controlled testing and UAT.",
        "Confirm workspace separation principles for Nursing Council and Medical Board records.",
        "Assign NDOH ICT owners for hosting, security review, backups, email, HTTPS/domain, and deployment planning.",
        "Agree on the production-readiness gates that must be completed before live use.",
        "Note the possible future DICT cloud assessment pathway after local/on-prem readiness evidence is prepared.",
    ]:
        document.add_paragraph(item, style="List Bullet")

    document.add_heading("Slide Notes", level=1)
    for idx, item in enumerate(SLIDES, 1):
        document.add_heading(f"Slide {idx}: {item['title']}", level=2)
        for bullet in item["bullets"]:
            document.add_paragraph(bullet, style="List Bullet")
        note = document.add_paragraph()
        note.add_run("Speaker note: ").bold = True
        note.add_run(item["notes"])

    document.add_heading("Suggested Demo Path", level=1)
    for item in [
        "Open the home page and show the public entry points.",
        "Open Nursing Council workspace and show registration/licensing/data quality views.",
        "Open Medical Board workspace and show separate office scope.",
        "Open Records Hub and explain controlled reference/registry management.",
        "Open documents, receipts, duplicate/missing-data review, and reports.",
        "Open mobile API readiness or Android local testing note.",
        "Close on production gates and ICT support required.",
    ]:
        document.add_paragraph(item, style="List Number")

    path = OUT / "TWG_Platform_Readiness_Speaker_Brief_20260609.docx"
    try:
        document.save(path)
    except PermissionError:
        path = OUT / "TWG_Platform_Readiness_Speaker_Brief_20260610_corrected.docx"
        document.save(path)
    return path


def build_pdf() -> Path:
    path = OUT / "TWG_Platform_Readiness_Presentation.pdf"
    dated_path = OUT / "TWG_Platform_Readiness_Presentation_20260609.pdf"
    doc = SimpleDocTemplate(
        str(path),
        pagesize=landscape((11 * inch, 6.1875 * inch)),
        rightMargin=0.45 * inch,
        leftMargin=0.45 * inch,
        topMargin=0.35 * inch,
        bottomMargin=0.35 * inch,
    )
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "SlideTitle",
        parent=styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=25,
        leading=29,
        textColor=colors.HexColor("#0A2B45"),
        spaceAfter=12,
    )
    bullet_style = ParagraphStyle(
        "SlideBullet",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=14,
        leading=18,
        leftIndent=16,
        firstLineIndent=-10,
        textColor=colors.HexColor("#0B1F2F"),
        spaceAfter=8,
    )
    def draw_page_label(canvas, doc_template):
        canvas.saveState()
        canvas.setFont("Helvetica-Bold", 9)
        canvas.setFillColor(colors.HexColor("#B07E24"))
        canvas.drawRightString(10.55 * inch, 5.86 * inch, ANNEX_LABEL)
        canvas.restoreState()

    story = []
    for idx, item in enumerate(SLIDES, 1):
        story.append(Paragraph(item["title"], title_style))
        if item.get("subtitle"):
            story.append(Paragraph(str(item["subtitle"]), styles["BodyText"]))
            story.append(Spacer(1, 0.12 * inch))
        for bullet in item["bullets"]:
            story.append(Paragraph(f"- {bullet}", bullet_style))
        image_path = screenshot(item.get("image") or "")
        if image_path:
            story.append(Spacer(1, 0.12 * inch))
            story.append(Image(str(image_path), width=4.6 * inch, height=2.45 * inch))
        if idx != len(SLIDES):
            story.append(PageBreak())
    doc.build(story, onFirstPage=draw_page_label, onLaterPages=draw_page_label)
    shutil.copyfile(path, dated_path)
    return path


def make_zip(paths: list[Path]) -> Path:
    zip_path = ROOT / "generated_reports" / "twg_presentation_20260609.zip"
    if zip_path.exists():
        zip_path.unlink()
    with ZipFile(zip_path, "w", ZIP_DEFLATED) as archive:
        for path in paths:
            archive.write(path, path.relative_to(OUT.parent))
    return zip_path


def write(path: Path, content: str) -> None:
    path.write_text(content.strip() + "\n", encoding="utf-8")


def main() -> None:
    reset_output()
    pptx_path = build_pptx()
    md_path = build_markdown()
    docx_path = build_docx()
    pdf_path = build_pdf()
    zip_path = make_zip([pptx_path, md_path, docx_path, pdf_path])
    print(f"PPTX={pptx_path}")
    print(f"DOCX={docx_path}")
    print(f"PDF={pdf_path}")
    print(f"MD={md_path}")
    print(f"ZIP={zip_path}")
    print(f"SLIDES={len(SLIDES)}")


if __name__ == "__main__":
    main()
